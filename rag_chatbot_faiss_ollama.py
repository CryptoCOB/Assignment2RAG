"""
ReqMind — Hybrid RAG (LangChain + FAISS + Ollama)
Student: Marc Harty (300818959)

Design choices:
- LangChain for PDF loading, splitting, embeddings, LLM orchestration
- FAISS for fast vector similarity search (via LangChain's FAISS wrapper)
- Ollama local models with lightweight routing (general/technical/vision)
"""
import os
import sys
from pathlib import Path
from typing import List

from dotenv import load_dotenv

from langchain_community.document_loaders import PyPDFLoader
from langchain_community.embeddings import SentenceTransformerEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document
# LangChain v1+ removed the old "chains" convenience import path for RetrievalQA in some builds.
# We implement a lightweight retrieval+stuff chain inline to avoid version coupling.
from langchain_ollama import ChatOllama
from langchain_text_splitters import RecursiveCharacterTextSplitter


QUIET = os.getenv("QUIET_MODE", "true").lower() == "true"

def _log(msg: str, level: str = "INFO", force: bool = False):
    """Minimal logging helper honoring QUIET flag.

    force=True overrides quiet (for errors / critical notices).
    """
    if not QUIET or force or level in {"ERROR", "WARN"}:
        print(msg)


class SimpleRetrievalQA:
    """Minimal RetrievalQA equivalent (chain_type='stuff').

    Provides an .invoke({"query": q}) interface returning
    { 'result': answer, 'source_documents': [Document, ...] }.
    """

    def __init__(self, llm: ChatOllama, retriever, k: int = 3, max_context_chars: int = 8000):
        self.llm = llm
        self.retriever = retriever
        self.k = k
        self.max_context_chars = max_context_chars

    def _build_context(self, docs):
        parts = []
        total = 0
        for i, d in enumerate(docs, 1):
            chunk = d.page_content.strip()
            if not chunk:
                continue
            if total + len(chunk) > self.max_context_chars:
                remaining = self.max_context_chars - total
                if remaining > 200:  # keep a reasonable tail fragment
                    chunk = chunk[:remaining]
                    parts.append(f"[Doc {i}]\n{chunk}")
                break
            parts.append(f"[Doc {i}]\n{chunk}")
            total += len(chunk)
        return "\n\n".join(parts)

    def invoke(self, inputs):
        if isinstance(inputs, dict):
            query = inputs.get("query") or inputs.get("question")
        else:
            query = str(inputs)
        if not query:
            raise ValueError("Query missing for RetrievalQA invocation")
        # Compatibility across LangChain versions
        if hasattr(self.retriever, "invoke"):
            docs = self.retriever.invoke(query)
        elif hasattr(self.retriever, "get_relevant_documents"):
            docs = self.retriever.get_relevant_documents(query)
        else:
            docs = self.retriever._get_relevant_documents(query, run_manager=None)
        context = self._build_context(docs)
        prompt = (
            "You are a helpful assistant answering questions about software system non-functional "
            "requirements. Use ONLY the provided context. If the answer is not in the context, say you "
            "cannot find it.\n\nCONTEXT:\n" + context + f"\n\nQUESTION: {query}\nANSWER:" 
        )
        response = self.llm.invoke(prompt)
        # ChatOllama returns an AIMessage or string depending on version
        answer = getattr(response, 'content', response)
        return {"result": answer, "source_documents": docs}

# Paths
ROOT = Path(__file__).parent
DATA_DIR = ROOT / "data"
CORPUS_DIR = Path(os.getenv("CORPUS_DIR", ROOT / "corpus")) if os.getenv("CORPUS_DIR") else (ROOT / "corpus")
FAISS_DIR = ROOT / "faiss_store" / "rag_300818959"


def load_env() -> None:
    load_dotenv(override=True)


def _load_pdf(file_path: Path) -> List:
    try:
        return PyPDFLoader(str(file_path)).load()
    except Exception as e:
        _log(f"[ERROR] PDF load failed: {file_path} -> {e}", level="ERROR", force=True)
        return []


def _load_md(file_path: Path) -> List:
    try:
        # Fallback: quick read as text; using TextLoader when available
        try:
            from langchain_community.document_loaders import TextLoader  # lazy import
            return TextLoader(str(file_path), autodetect_encoding=True).load()
        except Exception:
            text = file_path.read_text(encoding="utf-8", errors="ignore")
            return [Document(page_content=text, metadata={"source": str(file_path), "type": "markdown"})]
    except Exception as e:
        _log(f"[ERROR] MD load failed: {file_path} -> {e}", level="ERROR", force=True)
        return []


def _load_pptx(file_path: Path) -> List:
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        docs = []
        for idx, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = shape.text.strip()
                    if t:
                        texts.append(t)
            if not texts:
                continue
            content = "\n".join(texts)
            docs.append(
                Document(
                    page_content=content,
                    metadata={"source": str(file_path), "type": "pptx", "slide": idx},
                )
            )
        return docs
    except Exception as e:
        _log(f"[ERROR] PPTX load failed: {file_path} -> {e}", level="ERROR", force=True)
        return []


def _load_epub(file_path: Path) -> List:
    try:
        from ebooklib import epub
        from bs4 import BeautifulSoup
        book = epub.read_epub(str(file_path))
        docs = []
        for idx, item in enumerate(book.get_items_of_type(9), start=1):  # type 9 = DOCUMENT
            content_html = item.get_content().decode("utf-8", errors="ignore")
            soup = BeautifulSoup(content_html, "html.parser")
            text = soup.get_text(separator="\n", strip=True)
            if not text.strip():
                continue
            docs.append(
                Document(
                    page_content=text,
                    metadata={"source": str(file_path), "type": "epub", "chapter": idx},
                )
            )
        return docs
    except Exception as e:
        _log(f"[ERROR] EPUB load failed: {file_path} -> {e}", level="ERROR", force=True)
        return []


def load_documents() -> List:
    """Scan ./data and ./corpus for pdf, md, pptx, epub; return LangChain Documents."""
    all_docs: List = []
    # legacy data dir (pdfs per assignment spec)
    for pdf in sorted(DATA_DIR.glob("*.pdf")):
        all_docs.extend(_load_pdf(pdf))

    # corpus dir (pdf, md, pptx, epub)
    if not CORPUS_DIR.exists():
        CORPUS_DIR.mkdir(parents=True, exist_ok=True)
    corpus_files = list(CORPUS_DIR.glob("**/*"))
    for fp in corpus_files:
        if not fp.is_file():
            continue
        ext = fp.suffix.lower()
        if ext == ".pdf":
            all_docs.extend(_load_pdf(fp))
        elif ext in {".md", ".markdown", ".txt"}:
            all_docs.extend(_load_md(fp))
        elif ext == ".pptx":
            all_docs.extend(_load_pptx(fp))
        elif ext == ".epub":
            all_docs.extend(_load_epub(fp))

    if not all_docs:
        _log("[WARN] No documents found. Drop PDF/MD/PPTX/EPUB into ./corpus or ./data and re-run.", level="WARN")
    else:
        _log(f"[INFO] Loaded {len(all_docs)} documents from data/corpus")
    return all_docs


def split_documents(docs: List) -> List:
    chunk_size = int(os.getenv("CHUNK_SIZE", 1000))
    chunk_overlap = int(os.getenv("CHUNK_OVERLAP", 100))
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = splitter.split_documents(docs)
    _log(f"[INFO] Created {len(chunks)} chunks (chunk_size={chunk_size}, overlap={chunk_overlap})")
    if chunks and not QUIET:
        preview = chunks[0].page_content[:200].replace("\n", " ")
        _log(f"[PREVIEW] {preview}...")
    return chunks


def build_vectorstore(chunks: List) -> FAISS:
    embeddings = SentenceTransformerEmbeddings(model_name=os.getenv("EMBEDDING_MODEL", "all-MiniLM-L6-v2"))
    vs = FAISS.from_documents(chunks, embeddings)
    FAISS_DIR.parent.mkdir(parents=True, exist_ok=True)
    vs.save_local(str(FAISS_DIR))
    # touch a manifest timestamp
    (FAISS_DIR.parent / "INDEX_TIMESTAMP").write_text(str(max((d.metadata.get("source", "") for d in chunks), default="")))
    _log("[INFO] Vector store ready")
    return vs


def _index_mtime() -> float:
    if not FAISS_DIR.exists():
        return 0.0
    # Use directory mtime of FAISS_DIR (contains index files)
    try:
        return FAISS_DIR.stat().st_mtime
    except Exception:
        return 0.0


def _sources_newer_than_index(sources: List[Path]) -> bool:
    if not FAISS_DIR.exists():
        return True
    idx_mtime = _index_mtime()
    for fp in sources:
        try:
            if fp.stat().st_mtime > idx_mtime:
                return True
        except FileNotFoundError:
            return True
    return False


def load_vectorstore() -> FAISS:
    embeddings = SentenceTransformerEmbeddings(model_name=os.getenv("EMBEDDING_MODEL", "all-MiniLM-L6-v2"))
    return FAISS.load_local(str(FAISS_DIR), embeddings, allow_dangerous_deserialization=True)


def choose_llm(query: str) -> ChatOllama:
    q = query.lower()
    tech = any(k in q for k in ["code", "function", "method", "class", "implementation", "algorithm"])
    vision = any(k in q for k in ["diagram", "figure", "image", "table", "chart"])  # simple heuristic

    if vision:
        model = os.getenv("OLLAMA_VISION_MODEL", "llava-phi3:latest")
    elif tech:
        model = os.getenv("OLLAMA_TECHNICAL_MODEL", "deepseek-coder:6.7b")
    else:
        model = os.getenv("OLLAMA_PRIMARY_MODEL", "mistral:latest")

    return ChatOllama(model=model)


def build_chain(llm: ChatOllama, vectorstore: FAISS, k: int | None = None) -> SimpleRetrievalQA:
    if k is None:
        k = int(os.getenv("TOP_K_RETRIEVAL", 3))
    retriever = vectorstore.as_retriever(search_kwargs={"k": k})
    max_ctx = int(os.getenv("MAX_CONTEXT_CHARS", 8000))
    return SimpleRetrievalQA(llm=llm, retriever=retriever, k=k, max_context_chars=max_ctx)


def _summarize_chunks(llm: ChatOllama, texts: list[str], title: str = "Summary") -> str:
    """Summarize a list of text chunks and combine into a final concise summary."""
    partials: list[str] = []
    for idx, t in enumerate(texts, start=1):
        prompt = (
            f"You are summarizing part {idx} of a document titled '{title}'.\n"
            "Write a concise, faithful summary in bullet points. Do not invent information.\n\n"
            f"CONTENT:\n{t}\n\nSUMMARY:" 
        )
        resp = llm.invoke(prompt)
        partials.append(getattr(resp, "content", resp))
    # Combine partials
    combined = "\n\n".join(partials)
    final_prompt = (
        f"You are creating a final, unified summary for the document titled '{title}'.\n"
        "Synthesize the bullet points below into a clear, structured summary with headings for sections if present.\n\n"
        f"BULLET POINTS:\n{combined}\n\nFINAL SUMMARY:" 
    )
    final = llm.invoke(final_prompt)
    return getattr(final, "content", final)


def main():
    global QUIET
    load_env()
    first_name = os.getenv("STUDENT_NAME", "Marc").split()[0]
    _log(f"Starting {first_name} RAG Chatbot setup..")
    _log(f"[INFO] Corpus folder: {CORPUS_DIR}")

    # Prepare data
    docs = load_documents()
    if not docs:
        _log("[ERROR] No documents loaded. Drop PDF/MD/PPTX into ./corpus or ./data and re-run.", level="ERROR", force=True)
        sys.exit(1)
    chunks = split_documents(docs)
    if not chunks:
        _log("[ERROR] No chunks created. Check your source content.", level="ERROR", force=True)
        sys.exit(1)

    # Build / load FAISS with rebuild logic
    source_files: List[Path] = []
    source_files.extend(sorted(DATA_DIR.glob("*.pdf")))
    if CORPUS_DIR.exists():
        source_files.extend([p for p in CORPUS_DIR.glob("**/*") if p.is_file()])
    rebuild_if_newer = os.getenv("REBUILD_IF_NEWER", "true").lower() == "true"
    if (not FAISS_DIR.exists()) or (rebuild_if_newer and _sources_newer_than_index(source_files)):
        vectorstore = build_vectorstore(chunks)
    else:
        vectorstore = load_vectorstore()

    # Create an initial chain to validate end-to-end setup
    primary_llm = ChatOllama(model=os.getenv("OLLAMA_PRIMARY_MODEL", "mistral:latest"))
    _ = build_chain(primary_llm, vectorstore)
    _log("[INFO] RAG chain created")
    if QUIET:
        print("RAG ready. Type queries. ':help' for commands. (quiet mode)")
    else:
        print(f"RAG Chatbot {first_name} ready! Type 'exit' to quit.")
        print("[HELP] Commands: :files, :rebuild, :stats, :chunks, :test, :quiet, :verbose, :help")

    while True:
        try:
            q = input("\nYou: ").strip()
        except (EOFError, KeyboardInterrupt):
            print("\n[INFO] Exiting.")
            break
        if q.lower() == "exit":
            break
        if not q:
            continue

        # Commands
        if q.startswith(":"):
            cmd = q[1:].strip().lower()
            if cmd == "help":
                print("Commands:")
                print("  :files            - List source files")
                print("  :rebuild          - Force rebuild FAISS index")
                print("  :stats            - Show document and chunk statistics")
                print("  :chunks [filter]  - Preview chunks (optional: filter by filename)")
                print("  :test             - Test retrieval with sample queries")
                print("  :summary [filter] - Summarize all/filtered content")
                print("  :show [filter]    - Show first chunk matching filter")
                print("  :quiet            - Enable terse output")
                print("  :verbose          - Disable terse output")
                print("  exit              - Quit chatbot")
                continue
            if cmd == "quiet":
                QUIET = True
                print("Quiet mode ON")
                continue
            if cmd == "verbose":
                QUIET = False
                print("Quiet mode OFF")
                continue
            if cmd == "files":
                print("[Files]")
                for p in source_files:
                    try:
                        print(f"- {p.relative_to(ROOT)} ({p.stat().st_size/1024:.1f} KB)")
                    except Exception:
                        print(f"- {p}")
                continue
            if cmd == "rebuild":
                _log("[INFO] Rebuilding vector store...")
                vectorstore = build_vectorstore(chunks)
                _log("[INFO] Rebuild complete.")
                continue
            if cmd == "stats" or cmd.startswith("stats"):
                print("\n[Stats]")
                print(f"Total documents loaded: {len(docs)}")
                print(f"Total chunks created: {len(chunks)}")
                print(f"Chunk size: {os.getenv('CHUNK_SIZE', 1000)}")
                print(f"Chunk overlap: {os.getenv('CHUNK_OVERLAP', 100)}")
                print(f"Top-K retrieval: {os.getenv('TOP_K_RETRIEVAL', 3)}")
                print("Documents by type:")
                type_counts = {}
                for d in docs:
                    doc_type = d.metadata.get("type", "pdf")
                    type_counts[doc_type] = type_counts.get(doc_type, 0) + 1
                for t, c in sorted(type_counts.items()):
                    print(f"  {t}: {c}")
                print("Chunks by source:")
                source_counts = {}
                for c in chunks:
                    src = c.metadata.get("source", "unknown")
                    src_name = Path(src).name if src != "unknown" else "unknown"
                    source_counts[src_name] = source_counts.get(src_name, 0) + 1
                for s, cnt in sorted(source_counts.items()):
                    print(f"  {s}: {cnt} chunks")
                continue
            if cmd.startswith("chunks"):
                parts = cmd.split(maxsplit=1)
                flt = parts[1] if len(parts) > 1 else ""
                print(f"\n[Chunks] First 3{' filtered' if flt else ''}")
                shown = {}
                for idx, c in enumerate(chunks):
                    src = c.metadata.get("source", "unknown")
                    src_name = Path(src).name if src != "unknown" else "unknown"
                    if flt and flt.lower() not in src_name.lower():
                        continue
                    if shown.get(src_name, 0) >= 3:
                        continue
                    shown[src_name] = shown.get(src_name, 0) + 1
                    preview = c.page_content[:300].replace("\n", " ")
                    print(f"[{src_name} #{shown[src_name]}] {preview}...")
                continue
            if cmd == "test" or cmd.startswith("test"):
                print("\n[Test]")
                test_queries = [
                    "non-functional requirements",
                    "design patterns",
                    "assignment deliverables",
                ]
                for tq in test_queries:
                    print(f"Query: {tq}")
                    llm_test = ChatOllama(model=os.getenv("OLLAMA_PRIMARY_MODEL", "mistral:latest"))
                    qa_test = build_chain(llm_test, vectorstore)
                    result_test = qa_test.invoke({"query": tq})
                    sources_test = result_test.get("source_documents", [])
                    print(f"  Chunks: {len(sources_test)}")
                    for i, d in enumerate(sources_test, 1):
                        meta = d.metadata or {}
                        name = Path(meta.get("source", "?")).name
                        page = meta.get("page", meta.get("slide", meta.get("chapter", "?")))
                        snippet = d.page_content[:150].replace("\n", " ")
                        print(f"    {name}:{page} {snippet}...")
                continue
            if cmd.startswith("summary"):
                parts = cmd.split(maxsplit=1)
                flt = parts[1] if len(parts) > 1 else ""
                # collect all chunk texts matching filter (source path contains filter)
                target = [c for c in chunks if (flt.lower() in (c.metadata.get("source", "").lower())) or not flt]
                if not target:
                    print(f"[WARN] No chunks match filter '{flt}'.")
                    continue
                llm = ChatOllama(model=os.getenv("OLLAMA_PRIMARY_MODEL", "mistral:latest"))
                texts = [c.page_content[:3000] for c in target]  # safety cap per chunk
                summary = _summarize_chunks(llm, texts, title=f"Corpus Filter '{flt or 'ALL'}'")
                print("\n[Summary]\n" + summary)
                continue
            if cmd.startswith("show"):
                parts = cmd.split(maxsplit=1)
                flt = parts[1] if len(parts) > 1 else ""
                target = [c for c in chunks if (flt.lower() in (c.metadata.get("source", "").lower())) or not flt]
                if not target:
                    print(f"[WARN] No chunks match filter '{flt}'.")
                    continue
                print(f"[Show] First chunk of {len(target)} matched")
                first = target[0].page_content
                print(first[:2000])
                continue
            print("[WARN] Unknown command. Try :help")
            continue

        llm = choose_llm(q)
        qa = build_chain(llm, vectorstore)
        result = qa.invoke({"query": q})

        # Print sources
        sources = result.get("source_documents", [])
        if sources:
            if QUIET:
                compact = "; ".join(
                    f"{Path((d.metadata or {}).get('source','doc')).name}:{(d.metadata or {}).get('page', (d.metadata or {}).get('slide',(d.metadata or {}).get('chapter','?')))}" for d in sources
                )
                print(f"[Sources] {compact}")
            else:
                print("\n[Sources]")
                for i, d in enumerate(sources, 1):
                    meta = d.metadata or {}
                    name = meta.get("source") or meta.get("file_path") or meta.get("pdf") or "document"
                    page = meta.get("page", meta.get("page_number", meta.get("slide", "?")))
                    snippet = d.page_content[:160].replace("\n", " ")
                    print(f"{i}. {Path(name).name}:{page} {snippet}...")
        else:
            print("[Sources] none")

        answer = str(result.get("result") or result).strip()
        if QUIET:
            print(f"[Answer] {answer}")
        else:
            print("\n[Answer]\n" + answer)


if __name__ == "__main__":
    main()
